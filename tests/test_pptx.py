from __future__ import annotations

import base64
import os
from io import BytesIO

import pytest
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE
from pptx.util import Inches

from refill.filters import format_number_filter
from refill.pptx import PPTXFiller, PPTXParams, PPTXSpec, PPTXTemplate


@pytest.fixture
def dummy_b64_image():
    return (
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mNk+M/"
        + "wHwAEBgIApD5fRAAAAABJRU5ErkJggg=="
    )


@pytest.fixture
def dummy_data(dummy_b64_image):
    return {
        "item": {
            "id": 123,
            "image_url": "data:image/png;base64," + dummy_b64_image,
        },
        "stats": {
            "views": {"2021-11-01": 1200, "2021-11-02": 1400, "2021-11-03": 1100},
            "conversions": {"2021-11-01": 20, "2021-11-02": 40},
        },
    }


@pytest.fixture
def dummy_spec():
    return PPTXSpec(
        variables={"item_id": "item.id|str"},
        pictures={"item_image": "item.image_url|fetch"},
        tables={
            "*": {
                "keys": "stats.views|keys",
                "stubs": "=keys|selfie|format_date",
                "columns": {
                    "views": "stats.views|format_number",
                    "conversions": "stats.conversions|format_number",
                },
            }
        },
        charts={
            "*": {
                "keys": "stats.views|keys",
                "categories": "=keys|selfie",
                "series": {"views": "stats.views", "conversions": "stats.conversions"},
            }
        },
    )


@pytest.fixture
def dummy_template():
    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Item $item_id"

    picture_slide = prs.slides.add_slide(prs.slide_layouts[8])
    picture_slide.shapes.title.text = "Green is good!"
    for placeholder in picture_slide.placeholders:
        if placeholder.placeholder_format.type == PP_PLACEHOLDER_TYPE.PICTURE:  # type: ignore
            placeholder.name = "item_image"
            assert len(placeholder.element.getchildren()) < 3
            break

    table_slide = prs.slides.add_slide(prs.slide_layouts[5])
    table_slide.shapes.title.text = "Table for item ${item_id}"

    x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4)
    table_shape = table_slide.shapes.add_table(3, 3, x, y, cx, cy)
    table_shape.name = "stats_table"
    table = table_shape.table
    table.cell(0, 0).text = "Day"
    table.cell(0, 1).text = "Views"
    table.cell(0, 2).text = "Conversions"
    table.cell(1, 0).text = "$day1"
    table.cell(2, 0).text = "$day3"
    table.cell(1, 1).text = "n/a"
    table.cell(2, 1).text = "n/a"
    table.cell(1, 2).text = "n/a"
    table.cell(2, 2).text = "n/a"

    chart_slide = prs.slides.add_slide(prs.slide_layouts[5])
    chart_slide.shapes.title.text = "Chart for item ${item_id}"
    chart_data = ChartData()
    chart_data.categories = ["2021-01-01", "2021-01-02"]
    chart_data.add_series("views", (20, 25))
    chart_data.add_series("conversions", (2, 4))

    x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4)
    chart_shape = chart_slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, x, y, cx, cy, chart_data  # type: ignore
    )
    chart_shape.name = "stats_chart"

    buffer = BytesIO()
    prs.save(buffer)
    return PPTXTemplate(buffer)


def test_validate_spec(dummy_spec: PPTXSpec):
    PPTXParams.validate_spec(dummy_spec)


def test_apply_spec(dummy_spec: PPTXSpec, dummy_data, dummy_b64_image):
    assert dummy_spec.apply(dummy_data) == {
        "variables": {"item_id": str(dummy_data["item"]["id"])},
        "pictures": {"item_image": base64.b64decode(dummy_b64_image)},
        "tables": {
            "*": {
                "keys": ["2021-11-01", "2021-11-02", "2021-11-03"],
                "stubs": {
                    "2021-11-01": "Nov 1, 2021",
                    "2021-11-02": "Nov 2, 2021",
                    "2021-11-03": "Nov 3, 2021",
                },
                "columns": {
                    "views": {
                        "2021-11-01": "1,200",
                        "2021-11-02": "1,400",
                        "2021-11-03": "1,100",
                    },
                    "conversions": {"2021-11-01": "20", "2021-11-02": "40"},
                },
            },
        },
        "charts": {
            "*": {
                "keys": ["2021-11-01", "2021-11-02", "2021-11-03"],
                "categories": {
                    "2021-11-01": "2021-11-01",
                    "2021-11-02": "2021-11-02",
                    "2021-11-03": "2021-11-03",
                },
                "series": {
                    "views": {
                        "2021-11-01": 1200,
                        "2021-11-02": 1400,
                        "2021-11-03": 1100,
                    },
                    "conversions": {"2021-11-01": 20, "2021-11-02": 40},
                },
            }
        },
    }


def test_fill(dummy_spec: PPTXSpec, dummy_data, dummy_template: PPTXTemplate):
    buffer = BytesIO()
    PPTXFiller(dummy_spec, dummy_data).fill_to_file(dummy_template, buffer)
    buffer.seek(0)
    prs = Presentation(buffer)

    title_slide = prs.slides[0]
    assert title_slide.shapes.title.text == "Item " + str(dummy_data["item"]["id"])

    picture_slide = prs.slides[1]
    for placeholder in picture_slide.placeholders:
        if placeholder.placeholder_format.type == PP_PLACEHOLDER_TYPE.PICTURE:  # type: ignore
            assert placeholder.name == "item_image"
            assert len(placeholder.element.getchildren()) >= 3
            break

    table_slide = prs.slides[2]
    table_shape = table_slide.shapes[1]
    assert table_shape.name == "stats_table"
    assert table_shape.table.cell(1, 0).text == "Nov 1, 2021"
    assert table_shape.table.cell(2, 0).text == "Nov 3, 2021"
    assert table_shape.table.cell(1, 1).text == format_number_filter(
        dummy_data["stats"]["views"]["2021-11-01"], locale="en_US"
    )
    assert table_shape.table.cell(2, 1).text == format_number_filter(
        dummy_data["stats"]["views"]["2021-11-03"], locale="en_US"
    )
    assert table_shape.table.cell(1, 2).text == format_number_filter(
        dummy_data["stats"]["conversions"]["2021-11-01"], locale="en_US"
    )
    assert table_shape.table.cell(2, 2).text == "n/a"

    chart_slide = prs.slides[3]
    chart_shape = chart_slide.shapes[1]
    assert tuple(chart_shape.chart.plots[0].categories) == tuple(
        dummy_data["stats"]["views"].keys()
    )
    assert chart_shape.name == "stats_chart"
    assert chart_shape.chart.series[0].name == "views"
    assert chart_shape.chart.series[0].values == tuple(
        dummy_data["stats"]["views"].values()
    )
    assert chart_shape.chart.series[1].name == "conversions"
    assert chart_shape.chart.series[1].values[:-1] == tuple(
        dummy_data["stats"]["conversions"].values()
    )

    # write resulting pptx to file if environment variable is set
    output_pptx_file = os.environ.get("OUTPUT_PPTX_FILE")
    if output_pptx_file is not None:
        prs.save(output_pptx_file)
