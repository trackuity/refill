from __future__ import annotations

import re
import string
from abc import ABC, abstractmethod
from dataclasses import dataclass
from fnmatch import fnmatchcase
from io import BytesIO
from typing import IO, Dict, List, Optional, Union

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE
from typing_extensions import TypedDict

from .core import Filler, Params, Template
from .spec import Selector, Spec


class PPTXTableSpecDict(TypedDict):
    keys: Selector
    stubs: Selector
    columns: Dict[str, Selector]


class PPTXChartSpecDict(TypedDict):
    keys: Selector
    categories: Selector
    series: Dict[str, Selector]


@dataclass
class PPTXSpec(Spec):
    variables: Dict[str, Selector]
    pictures: Dict[str, Selector]
    tables: Dict[str, PPTXTableSpecDict]
    charts: Dict[str, PPTXChartSpecDict]


class PPTXTableParamsDict(TypedDict):
    keys: List[str]
    stubs: Dict[str, str]
    columns: Dict[str, Dict[str, str]]


class PPTXChartParamsDict(TypedDict):
    keys: List[str]
    categories: Dict[str, str]
    series: Dict[str, Dict[str, Union[int, float]]]


@dataclass
class PPTXParams(Params[PPTXSpec]):
    variables: Dict[str, str]
    pictures: Dict[str, bytes]
    tables: Dict[str, PPTXTableParamsDict]
    charts: Dict[str, PPTXChartParamsDict]


class PPTXTemplate(Template[PPTXParams]):
    def render_to_file(self, params: PPTXParams, file_object: IO[bytes]) -> None:
        self._params = params
        shape_substituters = self._create_shape_substituters()
        prs = Presentation(self._path_or_file)
        for slide in prs.slides:
            for shape in slide.shapes:
                for shape_substituter in shape_substituters:
                    shape_substituter.substitute_shape(shape)
        prs.save(file_object)

    def _create_shape_substituters(self) -> List[ShapeSubstituter]:
        return (
            [
                TextShapeSubstituter(self._params.variables),
                PicturePlaceholderSubstituter(self._params.pictures),
            ]
            + [
                TableShapeSubstituter(table_name, table_params)
                for (table_name, table_params) in self._params.tables.items()
            ]
            + [
                ChartShapeSubstituter(chart_name, chart_params)
                for (chart_name, chart_params) in self._params.charts.items()
            ]
        )


class PPTXFiller(Filler[PPTXSpec, PPTXParams, PPTXTemplate]):
    params_cls = PPTXParams


class ShapeSubstituter(ABC):
    @abstractmethod
    def substitute_shape(self, shape):
        ...


class TextShapeSubstituter(ShapeSubstituter):
    def __init__(self, variables: Dict[str, str]) -> None:
        super().__init__()
        self._variables = {
            name: (value if value is not None else "")
            for name, value in variables.items()
        }

    def substitute_shape(self, shape):
        if shape.has_text_frame:
            self.substitute_text_frame(shape.text_frame)

    def substitute_text_frame(self, text_frame):
        for paragraph in text_frame.paragraphs:
            if paragraph.runs:
                # Since powerpoint often splits text into multiple runs for some reason,
                # we combine the text from all runs, substitute that, and put the result
                # in the first run. The remaining runs are made empty. This implies that
                # the formatting from the first run will apply to everything in the end,
                # but templates can always use separate text frames if needed.
                first_run = paragraph.runs[0]
                first_run.text = self.substitute_text(
                    "".join(run.text for run in paragraph.runs)
                )
                for run in paragraph.runs[1:]:
                    run.text = ""

    def substitute_text(self, text: str) -> str:
        template = string.Template(text)
        return template.substitute(self._variables)


class PlaceholderSubstituter(ShapeSubstituter):
    def substitute_shape(self, shape):
        if shape.is_placeholder:
            self.substitute_placeholder(shape)

    @abstractmethod
    def substitute_placeholder(self, placeholder):
        ...


class PicturePlaceholderSubstituter(PlaceholderSubstituter):
    def __init__(self, pictures: Dict[str, bytes]) -> None:
        super().__init__()
        self._pictures = pictures

    def substitute_placeholder(self, placeholder):
        type_: PP_PLACEHOLDER_TYPE = placeholder.placeholder_format.type
        if type_ == PP_PLACEHOLDER_TYPE.PICTURE:  # type: ignore
            self.substitute_picture_placeholder(placeholder)

    def substitute_picture_placeholder(self, picture_placeholder):
        image = self._pictures.get(picture_placeholder.name)
        if image is not None:
            image_file = BytesIO(image)
            picture_placeholder.insert_picture(image_file)


class TableShapeSubstituter(ShapeSubstituter):
    def __init__(
        self, table_name_pattern: str, table_params: PPTXTableParamsDict
    ) -> None:
        super().__init__()
        self._table_name_pattern = table_name_pattern
        self._keys = table_params["keys"]
        self._stubs = table_params["stubs"]
        self._columns = table_params["columns"]

    def substitute_shape(self, shape):
        if shape.has_table and fnmatchcase(shape.name, self._table_name_pattern):
            self.substitute_table(shape.table)

    def substitute_table(self, table):
        column_index_values = []
        for i, row in enumerate(table.rows):
            if not column_index_values:  # first row is header
                for j, cell in enumerate(row.cells):
                    column_index_values.append(
                        self.derive_table_column_index_value(cell.text, j)
                    )
            else:
                row_index_value = None
                for j, cell in enumerate(row.cells):
                    if j == 0:
                        row_index_value = self.derive_table_row_index_value(
                            cell.text, i
                        )
                        for run in cell.text_frame.paragraphs[0].runs:
                            new_text = self.substitute_table_cell(
                                row_index_value, column_index_values[j], run.text
                            )
                            if new_text is not None:
                                run.text = new_text
                            break  # there should only be one run at most
                    else:
                        assert row_index_value is not None
                        for run in cell.text_frame.paragraphs[0].runs:
                            new_text = self.substitute_table_cell(
                                row_index_value, column_index_values[j], run.text
                            )
                            if new_text is not None:
                                run.text = new_text
                            break  # there should only be one run at most

    def derive_table_row_index_value(self, text: str, row_number: int) -> str:
        if text.startswith("$"):
            key_index = int(re.sub(r"\$[a-zA-Z_]*", "", text)) - 1
        else:
            key_index = row_number
        return self._keys[key_index]

    def derive_table_column_index_value(
        self, text: str, column_number: int
    ) -> Optional[str]:
        if column_number > 0:
            return text.lower()

    def substitute_table_cell(
        self,
        row_index_value: str,
        column_index_value: Optional[str],
        text: str,
    ) -> Optional[str]:
        if column_index_value is None:
            return self._stubs[row_index_value]
        else:
            return self._columns[column_index_value].get(row_index_value)


class ChartShapeSubstituter(ShapeSubstituter):
    def __init__(
        self, chart_name_pattern: str, chart_params: PPTXChartParamsDict
    ) -> None:
        super().__init__()
        self._chart_name_pattern = chart_name_pattern
        self._keys = chart_params["keys"]
        self._categories = chart_params["categories"]
        self._series = chart_params["series"]

    def substitute_shape(self, shape):
        if shape.has_chart and fnmatchcase(shape.name, self._chart_name_pattern):
            self.substitute_chart(shape.chart)

    def substitute_chart(self, chart):
        index_values = self.generate_chart_index_values(chart.plots[0].categories)
        metric_names = [self.derive_chart_metric_name(s.name) for s in chart.series]

        chart_data = ChartData()
        chart_data.categories = [
            self.get_chart_category(index_value) for index_value in index_values
        ]
        for metric_name in metric_names:
            values = [
                self.get_chart_value(index_value, metric_name)
                for index_value in index_values
            ]
            chart_data.add_series(metric_name, values)
        chart.replace_data(chart_data)

    def generate_chart_index_values(self, current_values) -> List[str]:
        return self._keys

    def derive_chart_metric_name(self, text: str) -> str:
        return text.lower()

    def get_chart_category(self, index_value: str) -> str:
        return self._categories[index_value]

    def get_chart_value(
        self, index_value: str, metric_name: str
    ) -> Optional[Union[int, float]]:
        return self._series[metric_name].get(index_value)
